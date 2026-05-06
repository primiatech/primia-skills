#!/usr/bin/env python3
"""
build_pptx.py — Constrói apresentação executiva em PowerPoint.

Uso:
    python build_pptx.py \
        --enriched /home/claude/output/concorrente1/enriched_ads.json ... \
        --analysis /home/claude/output/analysis.json \
        --output /mnt/user-data/outputs/apresentacao_concorrentes.pptx

Lembre-se: o Claude que usa esta skill deve ler /mnt/skills/public/pptx/SKILL.md
ANTES de rodar este script, para seguir as boas práticas oficiais.
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# Paleta
COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x78)
COLOR_ACCENT = RGBColor(0xE8, 0x7A, 0x22)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # background colorido
    left, top = Inches(0), Inches(0)
    width, height = prs.slide_width, prs.slide_height

    bg = slide.shapes.add_shape(1, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    # Title
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(9), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        tx2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(9), Inches(1))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(18)
        run2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    return slide


def add_section_header(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_ACCENT
    bg.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.7), Inches(3), Inches(9), Inches(1.5))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Faixa de título
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    return slide


def add_bullets(slide, bullets, top=Inches(1.2), left=Inches(0.5),
                width=Inches(9), height=Inches(5.5), font_size=18):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(8)


def add_text(slide, text, top, left, width, height, font_size=14, bold=False):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLOR_TEXT


def add_image_safe(slide, image_path, left, top, width, height):
    try:
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(str(image_path), left, top,
                                     width=width, height=height)
            return True
    except Exception:
        pass

    # placeholder
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_LIGHT
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    return False


def build(analysis: dict, all_ads: dict, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Capa
    title = analysis.get("title", "Análise Competitiva — Meta Ad Library")
    competitors = ", ".join(all_ads.keys())
    add_title_slide(prs, title, f"Concorrentes analisados: {competitors}")

    # 2. Sumário Executivo
    slide = add_content_slide(prs, "Sumário Executivo")
    summary = analysis.get("executive_summary", ["—"])
    add_bullets(slide, summary, font_size=20)

    # 3. Metodologia
    slide = add_content_slide(prs, "Metodologia")
    add_text(slide, analysis.get("methodology", "—"),
             top=Inches(1.2), left=Inches(0.5),
             width=Inches(12), height=Inches(5), font_size=14)

    # 4. Panorama
    slide = add_content_slide(prs, "Panorama Geral")
    panorama_text = analysis.get("panorama", "—")
    add_text(slide, panorama_text,
             top=Inches(1.2), left=Inches(0.5),
             width=Inches(12), height=Inches(2), font_size=14)

    # Tabela resumo
    panorama_data = []
    for competitor, info in all_ads.items():
        ads = info["ads"]
        active = sum(1 for a in ads if a.get("is_active"))
        panorama_data.append([competitor, str(len(ads)), str(active)])

    if panorama_data:
        rows, cols = len(panorama_data) + 1, 3
        tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(3.5),
                                           Inches(8), Inches(2.5))
        table = tbl_shape.table
        for i, h in enumerate(["Concorrente", "Total de anúncios", "Ativos"]):
            cell = table.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(14)
        for ri, row in enumerate(panorama_data, 1):
            for ci, val in enumerate(row):
                cell = table.cell(ri, ci)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(13)

    # 5. Por concorrente
    by_comp = analysis.get("by_competitor", {})
    for competitor, comp_data in by_comp.items():
        # Overview
        slide = add_content_slide(prs, f"{competitor} — Visão Geral")
        overview_lines = []
        if comp_data.get("perfil"):
            overview_lines.append(comp_data["perfil"])
        if comp_data.get("angulos"):
            overview_lines.append(f"Ângulos: {comp_data['angulos']}")
        if comp_data.get("funil"):
            overview_lines.append(f"Funil: {comp_data['funil']}")
        add_text(slide, "\n\n".join(overview_lines) or "—",
                 top=Inches(1.2), left=Inches(0.5),
                 width=Inches(12), height=Inches(5.5), font_size=15)

        # Top criativos (3 imagens lado a lado)
        top_ads = comp_data.get("top_ads", [])[:3]
        if top_ads:
            slide = add_content_slide(prs, f"{competitor} — Top 3 Anúncios Longevos")
            slot_w = Inches(4)
            slot_h = Inches(4.5)
            gap = Inches(0.3)
            start_left = Inches(0.5)
            top = Inches(1.3)

            for i, ad in enumerate(top_ads):
                left = start_left + i * (slot_w + gap)
                add_image_safe(slide, ad.get("screenshot"),
                               left, top, slot_w, slot_h)
                # caption
                caption = f"ID {ad.get('ad_id', '?')} — {ad.get('days_running', '?')} dias"
                add_text(slide, caption,
                         top=top + slot_h + Inches(0.1),
                         left=left, width=slot_w, height=Inches(0.4),
                         font_size=12, bold=True)
                if ad.get("destaque"):
                    add_text(slide, ad["destaque"][:120],
                             top=top + slot_h + Inches(0.5),
                             left=left, width=slot_w, height=Inches(1),
                             font_size=11)

    # 6. Comparativo
    matrix = analysis.get("comparative_matrix", [])
    if matrix:
        slide = add_content_slide(prs, "Análise Comparativa")
        headers = list(matrix[0].keys())
        rows = len(matrix) + 1
        cols = len(headers)
        tbl = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2),
                                     Inches(12.7), Inches(5.5)).table
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = str(h)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(11)
        for ri, row in enumerate(matrix, 1):
            for ci, h in enumerate(headers):
                cell = tbl.cell(ri, ci)
                cell.text = str(row.get(h, ""))
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(10)

    # 7. Oportunidades
    slide = add_content_slide(prs, "Oportunidades Identificadas")
    add_bullets(slide, analysis.get("opportunities", ["—"]), font_size=18)

    # 8. Recomendações
    recs = analysis.get("recommendations", [])
    # Divide em slides de até 5 recs cada
    chunk_size = 5
    for i in range(0, len(recs) or 1, chunk_size):
        chunk = recs[i:i + chunk_size]
        slide = add_content_slide(prs,
                                  f"Recomendações Estratégicas ({i + 1}–{i + len(chunk)})")
        bullets = [f"{j + 1 + i}. {r.get('titulo', '')}\n   → {r.get('evidencia', '')[:200]}"
                   for j, r in enumerate(chunk)]
        add_bullets(slide, bullets or ["—"], font_size=14)

    # 9. Encerramento
    add_section_header(prs, "Próximos Passos")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"[✓] Apresentação salva em {output}")


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--enriched", nargs="+", required=True)
    parser.add_argument("--analysis", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    all_ads = {}
    for path in args.enriched:
        with open(path, encoding="utf-8") as f:
            d = json.load(f)
        all_ads[d["metadata"]["advertiser"]] = {"ads": d.get("ads", [])}

    with open(args.analysis, encoding="utf-8") as f:
        analysis = json.load(f)

    build(analysis, all_ads, Path(args.output))


if __name__ == "__main__":
    main()
