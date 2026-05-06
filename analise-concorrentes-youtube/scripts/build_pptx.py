#!/usr/bin/env python3
"""
build_pptx.py

Gera apresentação executiva com top performers visuais e principais achados.

IMPORTANTE: Antes de rodar, leia /mnt/skills/public/pptx/SKILL.md.

Uso:
    python build_pptx.py --output-dir output/ --analysis-json analysis.json --output deck.pptx
"""
import argparse
import json
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Inches, Pt


# Paleta minimalista
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_PRIMARY = RGBColor(0x1F, 0x29, 0x37)
COLOR_ACCENT = RGBColor(0xEF, 0x44, 0x44)
COLOR_MUTED = RGBColor(0x6B, 0x72, 0x80)


def add_text_box(slide, x, y, w, h, text, size=14, bold=False, color=None, align=None):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_text_box(slide, 1.5, 6, 22, 3, title, size=40, bold=True, color=COLOR_PRIMARY)
    add_text_box(slide, 1.5, 10, 22, 1.5, subtitle, size=18, color=COLOR_MUTED)
    add_text_box(slide, 1.5, 17, 22, 1, datetime.now().strftime("%B de %Y"), size=12, color=COLOR_MUTED)
    return slide


def add_section_slide(prs, section_num, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 1.5, 7, 22, 1.5, section_num, size=14, color=COLOR_ACCENT, bold=True)
    add_text_box(slide, 1.5, 9, 22, 3, title, size=36, bold=True, color=COLOR_PRIMARY)
    return slide


def add_content_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text_box(slide, 1, 0.7, 22, 1.5, title, size=22, bold=True, color=COLOR_PRIMARY)
    return slide


def add_bullets(slide, bullets, x=1.5, y=2.5, w=22, max_h=14, size=16):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(max_h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = COLOR_PRIMARY


def add_image_safe(slide, image_path, x, y, w, h):
    try:
        if Path(image_path).exists():
            slide.shapes.add_picture(str(image_path), Cm(x), Cm(y), Cm(w), Cm(h))
            return True
    except Exception as e:
        print(f"  [erro imagem]: {e}", file=sys.stderr)
    return False


def build(prs, channels, top_videos_by_channel, analysis, output_dir):
    canais_str = ", ".join(c["title"] for c in channels)

    # 1. Capa
    add_title_slide(prs, "Análise Competitiva — YouTube", canais_str)

    # 2. Sumário Executivo
    slide = add_content_slide(prs, "Sumário Executivo")
    decisoes = analysis.get("decisoes_recomendadas", [])[:5]
    if decisoes:
        add_bullets(slide, decisoes, size=14)
    else:
        add_text_box(slide, 1.5, 4, 22, 2, "—", size=14)

    # 3. Metodologia
    slide = add_content_slide(prs, "Metodologia")
    meta = analysis.get("metodologia", {})
    bullets = [
        f"Modo: {meta.get('modo', 'Padrão')}",
        f"Janela: {meta.get('janela_dias', 180)} dias",
        f"Vídeos analisados: {meta.get('total_videos', '—')}",
        f"Transcrições: {meta.get('total_transcricoes', '—')}",
        f"Comentários: {meta.get('total_comentarios', '—')}",
    ]
    add_bullets(slide, bullets, size=14)

    # 4. Para cada canal: overview + top 3 thumbnails
    for ch in channels:
        slug = ch["slug"]
        s = analysis.get("channel_summaries", {}).get(slug, {})

        # Overview
        slide = add_content_slide(prs, ch["title"])
        bullets = [
            f"Inscritos: {ch.get('subscriber_count', 0):,}",
            f"Estágio funil dominante: {s.get('estagio_funil', '—')}",
            f"Ângulo dominante: {s.get('angulo_dominante', '—')}",
        ]
        if s.get("produto_ancora"):
            bullets.append(f"Produto âncora: {s['produto_ancora']}")
        if s.get("perfil"):
            bullets.append(s["perfil"])
        add_bullets(slide, bullets, size=14)

        # Top 3 thumbnails
        tops = top_videos_by_channel.get(slug, [])[:3]
        if tops:
            slide = add_content_slide(prs, f"{ch['title']} — Top vídeos")
            for i, v in enumerate(tops):
                x = 1.5 + i * 7.5
                thumb = output_dir / slug / "thumbnails" / f"{v['video_id']}.jpg"
                add_image_safe(slide, thumb, x=x, y=3, w=7, h=4)
                titulo_curto = v["title"][:60] + ("..." if len(v["title"]) > 60 else "")
                add_text_box(slide, x, 7.2, 7, 1.5, titulo_curto, size=10, bold=True)
                add_text_box(
                    slide, x, 8.5, 7, 1,
                    f"{v['view_count']:,} views | ratio {v['view_to_subscriber_ratio']}",
                    size=9, color=COLOR_MUTED,
                )

    # 5. Voz do Avatar — dores
    voz = analysis.get("voz_avatar_estruturada", {})
    dores = voz.get("dores", [])[:5]
    if dores:
        slide = add_content_slide(prs, "Voz do Avatar — Dores principais")
        bullets = []
        for d in dores:
            if isinstance(d, dict):
                titulo = d.get("titulo") or d.get("padrao", "")
                quotes = d.get("quotes", [])
                bullet = titulo
                if quotes:
                    bullet += f'  "{quotes[0][:120]}"'
                bullets.append(bullet)
            else:
                bullets.append(str(d))
        add_bullets(slide, bullets, size=12)

    # 6. Voz do Avatar — perguntas (= ideias de conteúdo)
    perguntas = voz.get("perguntas", [])[:5]
    if perguntas:
        slide = add_content_slide(prs, "Perguntas sem resposta = Ideias de conteúdo")
        bullets = []
        for p in perguntas:
            if isinstance(p, dict):
                bullets.append(p.get("titulo") or p.get("padrao", ""))
            else:
                bullets.append(str(p))
        add_bullets(slide, bullets, size=14)

    # 7. Comparativo
    matriz = analysis.get("comparativo_matrix", [])
    if matriz:
        slide = add_content_slide(prs, "Análise Comparativa")
        # Renderiza como tabela
        headers = list(matriz[0].keys())
        rows = len(matriz) + 1
        cols = len(headers)
        table = slide.shapes.add_table(
            rows, cols, Cm(1), Cm(2.5), Cm(22), Cm(min(rows * 0.8, 12))
        ).table
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(10)
        for ri, row in enumerate(matriz, 1):
            for ci, h in enumerate(headers):
                cell = table.cell(ri, ci)
                cell.text = str(row.get(h, ""))
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(9)

    # 8. Gaps e oportunidades
    gaps = analysis.get("gaps", [])
    if gaps:
        slide = add_content_slide(prs, "Gaps e Oportunidades")
        bullets = []
        for g in gaps[:6]:
            if isinstance(g, dict):
                bullets.append(f"{g.get('tema', '')} — {g.get('justificativa', '')}")
            else:
                bullets.append(str(g))
        add_bullets(slide, bullets, size=13)

    # 9. Banco de hooks
    hooks = analysis.get("banco_de_hooks", [])[:6]
    if hooks:
        slide = add_content_slide(prs, "Banco de Hooks")
        bullets = []
        for h in hooks:
            if isinstance(h, dict):
                bullets.append(f"{h.get('tipo', '')}: {h.get('exemplo', '')}")
            else:
                bullets.append(str(h))
        add_bullets(slide, bullets, size=13)

    # 10. Recomendações (até 3 slides, 3 recs por slide)
    recs = analysis.get("recomendacoes", [])
    chunks = [recs[i : i + 3] for i in range(0, len(recs), 3)]
    for i, chunk in enumerate(chunks[:3], 1):
        slide = add_content_slide(prs, f"Recomendações ({i}/{len(chunks[:3])})")
        bullets = []
        for r in chunk:
            o_que = r.get("o_que", "")
            como = r.get("como", "")
            bullets.append(f"{o_que}  →  {como[:120]}")
        add_bullets(slide, bullets, size=12)

    # 11. Próximos passos
    slide = add_content_slide(prs, "Próximos Passos")
    proximos = analysis.get("proximos_passos", [
        "Priorizar 1-2 recomendações para teste imediato",
        "Definir métricas de validação (view ratio, comentários, conversão de lead)",
        "Reanalisar após 30-60 dias para medir efeito",
    ])
    add_bullets(slide, proximos, size=14)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--analysis-json", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    base = Path(args.output_dir)

    channels_file = base / "channels.json"
    if not channels_file.exists():
        channels_file = base.parent / "channels.json"
    with open(channels_file, "r", encoding="utf-8") as f:
        channels = json.load(f)

    top_videos_by_channel = {}
    for ch in channels:
        slug = ch["slug"]
        top_file = base / slug / "top_videos.json"
        if top_file.exists():
            with open(top_file, "r", encoding="utf-8") as f:
                tops = json.load(f)
            for v in tops:
                v["channel_title"] = ch["title"]
            top_videos_by_channel[slug] = tops

    with open(args.analysis_json, "r", encoding="utf-8") as f:
        analysis = json.load(f)

    # Wide format (16:9)
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    build(prs, channels, top_videos_by_channel, analysis, base)

    Path(args.output).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"Apresentação salva em {args.output}", file=sys.stderr)


if __name__ == "__main__":
    main()
